import os
import json
import uuid
from datetime import datetime
from typing import Dict, List, Any, Optional
from pptx import Presentation as PPTXPresentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import requests
from PIL import Image
import io

from src.models.presentation import db, Presentation, Slide, Citation
from src.services.ai_service import AIService
from src.services.image_service import ImageService

class PresentationService:
    """Service for generating PowerPoint presentations using python-pptx."""
    
    def __init__(self):
        self.ai_service = AIService()
        self.image_service = ImageService()
        self.themes = {
            'corporate': {
                'primary_color': RGBColor(31, 41, 55),    # #1f2937
                'secondary_color': RGBColor(59, 130, 246), # #3b82f6
                'accent_color': RGBColor(16, 185, 129),    # #10b981
                'background_color': RGBColor(255, 255, 255),
                'text_color': RGBColor(31, 41, 55),
                'font_name': 'Calibri'
            },
            'startup': {
                'primary_color': RGBColor(124, 58, 237),   # #7c3aed
                'secondary_color': RGBColor(245, 158, 11), # #f59e0b
                'accent_color': RGBColor(239, 68, 68),     # #ef4444
                'background_color': RGBColor(255, 255, 255),
                'text_color': RGBColor(55, 65, 81),
                'font_name': 'Calibri'
            },
            'academic': {
                'primary_color': RGBColor(55, 65, 81),     # #374151
                'secondary_color': RGBColor(99, 102, 241), # #6366f1
                'accent_color': RGBColor(5, 150, 105),     # #059669
                'background_color': RGBColor(255, 255, 255),
                'text_color': RGBColor(55, 65, 81),
                'font_name': 'Calibri'
            }
        }
    
    def generate_presentation(self, presentation_id: int, research_data: List[Dict[str, Any]] = None) -> bool:
        """
        Generate a complete PowerPoint presentation.
        
        Args:
            presentation_id: Database ID of the presentation to generate
            research_data: Optional research data to use for content
            
        Returns:
            True if generation successful, False otherwise
        """
        try:
            # Get presentation from database
            presentation = Presentation.query.get(presentation_id)
            if not presentation:
                return False
            
            # Update status
            presentation.status = 'planning'
            presentation.progress = 20
            presentation.current_step = 'Creating presentation structure'
            db.session.commit()
            
            # Analyze prompt and create outline
            analysis = self.ai_service.analyze_prompt(presentation.original_prompt)
            outline = self.ai_service.generate_presentation_outline(
                presentation.original_prompt,
                presentation.slide_count,
                analysis
            )
            
            # Update status
            presentation.status = 'generating'
            presentation.progress = 40
            presentation.current_step = 'Generating slide content'
            db.session.commit()
            
            # Create PPTX presentation
            pptx = PPTXPresentation()
            theme = self.themes.get(presentation.theme, self.themes['corporate'])
            
            # Generate slides
            citations = []
            citation_counter = 1
            
            for i, slide_info in enumerate(outline.get('slides', [])):
                # Generate content for this slide
                slide_research_data = research_data[:3] if research_data else []
                slide_content = self.ai_service.generate_slide_content(
                    slide_info,
                    slide_research_data,
                    {'title': presentation.title}
                )
                
                # Get images for this slide
                slide_images = self.image_service.get_images_for_slide(slide_content, slide_info['type'])
                
                # Create slide in PPTX
                slide = self._create_slide(pptx, slide_info, slide_content, theme)
                
                # Add images to slide if available
                if slide_images:
                    self._add_images_to_slide(slide, slide_images[:1])  # Add max 1 image per slide
                
                # Save slide to database
                db_slide = Slide(
                    presentation_id=presentation.id,
                    slide_number=slide_info['slide_number'],
                    title=slide_content.get('title', slide_info['title']),
                    slide_type=slide_info['type'],
                    content_json=json.dumps(slide_content),
                    speaker_notes=slide_content.get('speaker_notes', ''),
                    image_urls=json.dumps([img['url'] for img in slide_images]),
                    chart_data=json.dumps({})   # TODO: Add chart data
                )
                db.session.add(db_slide)
                
                # Handle citations
                for citation_id in slide_content.get('citations', []):
                    if citation_id <= len(slide_research_data):
                        source = slide_research_data[citation_id - 1]
                        citation = Citation(
                            presentation_id=presentation.id,
                            citation_number=citation_counter,
                            source_type='web',
                            title=source.get('title', 'Unknown Source'),
                            url=source.get('url', ''),
                            accessed_date=datetime.utcnow().date()
                        )
                        db.session.add(citation)
                        citations.append(citation)
                        citation_counter += 1
                
                # Update progress
                progress = 40 + int((i + 1) / len(outline['slides']) * 40)
                presentation.progress = progress
                db.session.commit()
            
            # Add references slide if citations exist
            if citations:
                self._create_references_slide(pptx, citations, theme)
            
            # Update status
            presentation.status = 'assembling'
            presentation.progress = 90
            presentation.current_step = 'Finalizing presentation file'
            db.session.commit()
            
            # Save PPTX file
            file_path = self._save_presentation(pptx, presentation)
            
            # Update presentation record
            presentation.status = 'completed'
            presentation.progress = 100
            presentation.current_step = 'Completed'
            presentation.file_path = file_path
            presentation.file_size = os.path.getsize(file_path) if os.path.exists(file_path) else 0
            presentation.completed_at = datetime.utcnow()
            db.session.commit()
            
            return True
            
        except Exception as e:
            print(f"Error generating presentation: {e}")
            # Update presentation with error
            presentation.status = 'failed'
            presentation.error_message = str(e)
            db.session.commit()
            return False
    
    def _add_images_to_slide(self, slide: Any, images: List[Dict[str, Any]]) -> None:
        """Add images to a slide."""
        try:
            for i, image_info in enumerate(images):
                local_path = image_info.get('local_path')
                if local_path and os.path.exists(local_path):
                    # Position image on the right side of the slide
                    left = Inches(7)  # Right side
                    top = Inches(2 + i * 2)  # Stacked vertically
                    width = Inches(3)
                    
                    try:
                        slide.shapes.add_picture(local_path, left, top, width=width)
                    except Exception as e:
                        print(f"Error adding image to slide: {e}")
                        continue
        except Exception as e:
            print(f"Error in _add_images_to_slide: {e}")
    
    def _create_slide(self, pptx: PPTXPresentation, slide_info: Dict[str, Any], slide_content: Dict[str, Any], theme: Dict[str, Any]) -> Any:
        """Create a single slide in the presentation."""
        slide_type = slide_info.get('type', 'content')
        
        if slide_type == 'title':
            return self._create_title_slide(pptx, slide_content, theme)
        elif slide_type == 'agenda':
            return self._create_agenda_slide(pptx, slide_content, theme)
        elif slide_type == 'content':
            return self._create_content_slide(pptx, slide_content, theme)
        elif slide_type == 'conclusion':
            return self._create_conclusion_slide(pptx, slide_content, theme)
        else:
            return self._create_content_slide(pptx, slide_content, theme)
    
    def _create_title_slide(self, pptx: PPTXPresentation, content: Dict[str, Any], theme: Dict[str, Any]) -> Any:
        """Create a title slide."""
        slide_layout = pptx.slide_layouts[0]  # Title slide layout
        slide = pptx.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = content.get('title', 'Presentation Title')
        title.text_frame.paragraphs[0].font.name = theme['font_name']
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = theme['primary_color']
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Set subtitle if available
        if slide.placeholders.count > 1:
            subtitle = slide.placeholders[1]
            subtitle.text = content.get('subtitle', 'Generated by Prompt2Presentation')
            subtitle.text_frame.paragraphs[0].font.name = theme['font_name']
            subtitle.text_frame.paragraphs[0].font.size = Pt(24)
            subtitle.text_frame.paragraphs[0].font.color.rgb = theme['secondary_color']
            subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        return slide
    
    def _create_agenda_slide(self, pptx: PPTXPresentation, content: Dict[str, Any], theme: Dict[str, Any]) -> Any:
        """Create an agenda slide."""
        slide_layout = pptx.slide_layouts[1]  # Title and content layout
        slide = pptx.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = content.get('title', 'Agenda')
        title.text_frame.paragraphs[0].font.name = theme['font_name']
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = theme['primary_color']
        
        # Add agenda items
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        bullet_points = content.get('bullet_points', content.get('content', {}).get('items', []))
        for i, item in enumerate(bullet_points):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = item
            p.font.name = theme['font_name']
            p.font.size = Pt(24)
            p.font.color.rgb = theme['text_color']
            p.level = 0
        
        return slide
    
    def _create_content_slide(self, pptx: PPTXPresentation, content: Dict[str, Any], theme: Dict[str, Any]) -> Any:
        """Create a content slide with bullet points."""
        slide_layout = pptx.slide_layouts[1]  # Title and content layout
        slide = pptx.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = content.get('title', 'Content Slide')
        title.text_frame.paragraphs[0].font.name = theme['font_name']
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = theme['primary_color']
        
        # Add content
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        bullet_points = content.get('bullet_points', [])
        for i, point in enumerate(bullet_points):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = point
            p.font.name = theme['font_name']
            p.font.size = Pt(20)
            p.font.color.rgb = theme['text_color']
            p.level = 0
        
        # Add speaker notes
        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        notes_text_frame.text = content.get('speaker_notes', '')
        
        return slide
    
    def _create_conclusion_slide(self, pptx: PPTXPresentation, content: Dict[str, Any], theme: Dict[str, Any]) -> Any:
        """Create a conclusion slide."""
        slide_layout = pptx.slide_layouts[1]  # Title and content layout
        slide = pptx.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = content.get('title', 'Conclusion')
        title.text_frame.paragraphs[0].font.name = theme['font_name']
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = theme['primary_color']
        
        # Add conclusion points
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        bullet_points = content.get('bullet_points', ['Thank you for your attention', 'Questions?'])
        for i, point in enumerate(bullet_points):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = point
            p.font.name = theme['font_name']
            p.font.size = Pt(24)
            p.font.color.rgb = theme['text_color']
            p.level = 0
        
        return slide
    
    def _create_references_slide(self, pptx: PPTXPresentation, citations: List[Citation], theme: Dict[str, Any]) -> Any:
        """Create a references slide."""
        slide_layout = pptx.slide_layouts[1]  # Title and content layout
        slide = pptx.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = "References"
        title.text_frame.paragraphs[0].font.name = theme['font_name']
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = theme['primary_color']
        
        # Add references
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        for i, citation in enumerate(citations):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = f"{citation.citation_number}. {citation.title} - {citation.url}"
            p.font.name = theme['font_name']
            p.font.size = Pt(14)
            p.font.color.rgb = theme['text_color']
            p.level = 0
        
        return slide
    
    def _save_presentation(self, pptx: PPTXPresentation, presentation: Presentation) -> str:
        """Save the PPTX presentation to file."""
        # Create presentations directory if it doesn't exist
        presentations_dir = os.path.join(os.path.dirname(__file__), '..', '..', 'presentations')
        os.makedirs(presentations_dir, exist_ok=True)
        
        # Generate filename
        safe_title = "".join(c for c in presentation.title if c.isalnum() or c in (' ', '-', '_')).rstrip()
        filename = f"{safe_title}_{presentation.session_id[:8]}.pptx"
        file_path = os.path.join(presentations_dir, filename)
        
        # Save the presentation
        pptx.save(file_path)
        
        return file_path
    
    def add_image_to_slide(self, slide: Any, image_url: str, left: float = 6.0, top: float = 2.0, width: float = 4.0) -> bool:
        """
        Add an image to a slide.
        
        Args:
            slide: The slide object
            image_url: URL or path to the image
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            
        Returns:
            True if image added successfully, False otherwise
        """
        try:
            if image_url.startswith('http'):
                # Download image from URL
                response = requests.get(image_url, timeout=10)
                response.raise_for_status()
                image_stream = io.BytesIO(response.content)
            else:
                # Local file
                image_stream = image_url
            
            # Add image to slide
            slide.shapes.add_picture(
                image_stream,
                Inches(left),
                Inches(top),
                width=Inches(width)
            )
            
            return True
            
        except Exception as e:
            print(f"Error adding image to slide: {e}")
            return False
    
    def create_chart_slide(self, pptx: PPTXPresentation, title: str, chart_data: Dict[str, Any], theme: Dict[str, Any]) -> Any:
        """Create a slide with a chart."""
        slide_layout = pptx.slide_layouts[5]  # Title and content layout
        slide = pptx.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.name = theme['font_name']
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme['primary_color']
        
        # Create chart data
        chart_data_obj = CategoryChartData()
        chart_data_obj.categories = chart_data.get('categories', ['Category 1', 'Category 2', 'Category 3'])
        chart_data_obj.add_series('Series 1', chart_data.get('values', [10, 20, 30]))
        
        # Add chart
        x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data_obj
        ).chart
        
        return slide

